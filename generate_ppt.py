"""
INUBIL Verify — Présentation soutenance
python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1a, 0x1a, 0x2e)
NAVY_MED  = RGBColor(0x16, 0x21, 0x3e)
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_DARK = RGBColor(0x9A, 0x7A, 0x2A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xF4, 0xF4, 0xF8)
DGRAY     = RGBColor(0x44, 0x44, 0x55)
GREEN     = RGBColor(0x2D, 0x7A, 0x3A)
RED       = RGBColor(0xB9, 0x1C, 0x1C)
ORANGE    = RGBColor(0xC2, 0x71, 0x0C)
PURPLE    = RGBColor(0x7E, 0x57, 0xC2)
BLUE      = RGBColor(0x22, 0x5C, 0x8A)

# ── Dimensions 16:9 ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Chemins ressources ────────────────────────────────────────────────────────
ASSETS   = r"backend\src\assets"
ICONS    = os.path.join(ASSETS, "icons")
TECHLOGO = os.path.join(ASSETS, "tech-logos")
INUBIL_LOGO = os.path.join(ASSETS, "inubil-logo.png")
EVRY_LOGO   = os.path.join(ASSETS, "evry-logo.png")

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, color, x, y, w, h):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txb(slide, text, x, y, w, h, size=16, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb

def img(slide, path, x, y, h=None, w=None):
    if not os.path.exists(path):
        return None
    try:
        if h and w:
            return slide.shapes.add_picture(path, x, y, width=w, height=h)
        elif h:
            return slide.shapes.add_picture(path, x, y, height=h)
        else:
            return slide.shapes.add_picture(path, x, y)
    except Exception:
        return None

def icon(slide, name, x, y, size=Inches(0.55)):
    return img(slide, os.path.join(ICONS, name + ".png"), x, y, h=size)

def tech(slide, name, x, y, size=Inches(0.65)):
    return img(slide, os.path.join(TECHLOGO, name + ".png"), x, y, h=size)

def num_badge(slide, n, total=20):
    txb(slide, f"{n} / {total}",
        W - Inches(1.3), H - Inches(0.4),
        Inches(1.1), Inches(0.32),
        size=9, color=GOLD, align=PP_ALIGN.RIGHT)

def section_pill(slide, text):
    rect(slide, GOLD, Inches(0.35), Inches(0.22), Inches(2.6), Inches(0.34))
    txb(slide, text.upper(), Inches(0.43), Inches(0.23), Inches(2.5), Inches(0.3),
        size=9, color=NAVY, bold=True)

def h_rule(slide, x, y, w, color=GOLD, th=Inches(0.04)):
    rect(slide, color, x, y, w, th)

def kpi(slide, label, val, unit, x, y, bg_c=NAVY_MED):
    rect(slide, bg_c, x, y, Inches(2.85), Inches(1.55))
    h_rule(slide, x, y + Inches(1.35), Inches(2.85), GOLD, Inches(0.05))
    txb(slide, val, x, y + Inches(0.08), Inches(2.85), Inches(0.75),
        size=38, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    txb(slide, unit, x, y + Inches(0.76), Inches(2.85), Inches(0.35),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)
    txb(slide, label, x, y + Inches(1.08), Inches(2.85), Inches(0.32),
        size=11, color=LGRAY, align=PP_ALIGN.CENTER)

def notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def icon_card(slide, icon_name, title, desc, accent, x, y, w=Inches(6.1), h=Inches(1.35)):
    rect(slide, WHITE, x, y, w, h)
    rect(slide, accent, x, y, Inches(0.06), h)
    icon(slide, icon_name, x + Inches(0.18), y + Inches(0.2), size=Inches(0.85))
    txb(slide, title, x + Inches(1.2), y + Inches(0.1), w - Inches(1.3), Inches(0.42),
        size=14, color=NAVY, bold=True)
    txb(slide, desc, x + Inches(1.2), y + Inches(0.52), w - Inches(1.3), Inches(0.75),
        size=11, color=DGRAY, wrap=True)

def step_box(slide, num, title, desc, x, y, active=False):
    bg_c = NAVY if active else LGRAY
    fg_c = WHITE if active else NAVY
    sm_c = GOLD if active else DGRAY
    rect(slide, bg_c, x, y, Inches(2.35), Inches(2.0))
    txb(slide, num, x + Inches(0.12), y + Inches(0.1), Inches(0.5), Inches(0.5),
        size=22, color=GOLD, bold=True)
    txb(slide, title, x + Inches(0.1), y + Inches(0.58), Inches(2.1), Inches(0.5),
        size=13, color=fg_c, bold=True)
    txb(slide, desc, x + Inches(0.1), y + Inches(1.08), Inches(2.1), Inches(0.85),
        size=11, color=sm_c, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)
rect(s, NAVY_MED, Inches(9.2), Inches(5.2), Inches(4.13), Inches(2.3))
rect(s, GOLD, Inches(9.2), Inches(5.2), Inches(4.13), Inches(0.05))

img(s, INUBIL_LOGO, Inches(0.55), Inches(0.25), h=Inches(1.5))
rect(s, GOLD, Inches(2.55), Inches(0.35), Inches(0.04), Inches(1.1))
img(s, EVRY_LOGO, Inches(2.75), Inches(0.3), h=Inches(2.5))

rect(s, GOLD_DARK, Inches(0.55), Inches(2.0), Inches(4.5), Inches(0.38))
txb(s, "SOUTENANCE DE FIN D'ETUDES  -  JUIN 2026",
    Inches(0.65), Inches(2.02), Inches(4.3), Inches(0.35),
    size=10, color=WHITE, bold=True)

txb(s, "INUBIL VERIFY",
    Inches(0.55), Inches(2.55), Inches(11), Inches(1.15),
    size=56, color=GOLD, bold=True)

txb(s, "Plateforme Blockchain de Certification\net d'Authentification de Diplomes",
    Inches(0.55), Inches(3.75), Inches(9), Inches(1.1), size=21, color=WHITE)

h_rule(s, Inches(0.55), Inches(4.95), Inches(12.4))

rect(s, NAVY_MED, Inches(0.55), Inches(5.1), Inches(5.9), Inches(1.25))
rect(s, GOLD, Inches(0.55), Inches(5.1), Inches(0.06), Inches(1.25))
txb(s, "TCHENTCHEU JIAGAM FLANC BEL",
    Inches(0.78), Inches(5.18), Inches(5.6), Inches(0.42), size=14, color=WHITE, bold=True)
txb(s, "Developpeur Backend", Inches(0.78), Inches(5.6), Inches(5.6), Inches(0.35),
    size=12, color=GOLD)
txb(s, "NestJS  -  PostgreSQL  -  Blockchain Polygon",
    Inches(0.78), Inches(5.95), Inches(5.6), Inches(0.32), size=10, color=LGRAY, italic=True)

rect(s, NAVY_MED, Inches(6.7), Inches(5.1), Inches(6.25), Inches(1.25))
rect(s, GOLD, Inches(6.7), Inches(5.1), Inches(0.06), Inches(1.25))
txb(s, "NGANGUE TSAFACK BELVIE SCINDIE",
    Inches(6.93), Inches(5.18), Inches(5.95), Inches(0.42), size=14, color=WHITE, bold=True)
txb(s, "Developpeuse Frontend & Maquettes",
    Inches(6.93), Inches(5.6), Inches(5.95), Inches(0.35), size=12, color=GOLD)
txb(s, "Angular  -  UI/UX  -  Figma",
    Inches(6.93), Inches(5.95), Inches(5.95), Inches(0.32), size=10, color=LGRAY, italic=True)

txb(s, "ISTAMA INUBIL  -  Institut Universitaire Bilingue du Littoral  -  Douala, Cameroun",
    Inches(0.55), Inches(6.6), Inches(12.4), Inches(0.38),
    size=11, color=LGRAY, align=PP_ALIGN.CENTER)

notes(s, "Bonjour a tous. Je suis TCHENTCHEU JIAGAM FLANC BEL et voici NGANGUE TSAFACK BELVIE SCINDIE. "
         "Nous vous presentons INUBIL Verify, notre projet de fin d'etudes developpe pour ISTAMA INUBIL. "
         "INUBIL Verify est une plateforme web qui permet d'emettre des diplomes certifies sur la blockchain "
         "et de les verifier en quelques secondes, depuis n'importe ou dans le monde. "
         "Je vais vous presenter la partie backend et architecture, puis NGANGUE vous montrera le frontend lors de la demonstration.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

txb(s, "Plan de la presentation",
    Inches(0.5), Inches(0.3), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

plan = [
    ("01", "Problematique",       "Verification de diplomes : un processus sans solution",    "alert"),
    ("02", "Solution",            "INUBIL Verify en un coup d'oeil",                          "zap"),
    ("03", "Architecture",        "Comment les composants s'articulent",                       "layers"),
    ("04", "Stack technologique", "Choix techniques et justifications",                        "cpu"),
    ("05", "Fonctionnalites",     "Emission, verification, securite",                          "award"),
    ("06", "Infrastructure",      "Docker, Cloudflare R2, Polygon Amoy",                      "server"),
    ("07", "Demo & Bilan",        "Resultats, defis, perspectives",                           "trending-up"),
]
for i, (num, title, desc, ico) in enumerate(plan):
    row = i // 2
    col = i % 2
    x = Inches(0.42 + col * 6.45)
    y = Inches(1.45 + row * 1.7)
    if i == 6:
        x, y = Inches(0.42), Inches(1.45 + 3 * 1.7)
    rect(s, WHITE, x, y, Inches(6.1), Inches(1.45))
    rect(s, NAVY, x, y, Inches(0.06), Inches(1.45))
    rect(s, NAVY, x, y, Inches(6.1), Inches(0.55))
    txb(s, num, x + Inches(0.15), y + Inches(0.09), Inches(0.5), Inches(0.38),
        size=16, color=GOLD, bold=True)
    txb(s, title, x + Inches(0.72), y + Inches(0.1), Inches(5.2), Inches(0.38),
        size=14, color=WHITE, bold=True)
    icon(s, ico, x + Inches(0.15), y + Inches(0.65), size=Inches(0.65))
    txb(s, desc, x + Inches(1.0), y + Inches(0.7), Inches(4.9), Inches(0.65),
        size=12, color=DGRAY, italic=True)

num_badge(s, 2)
notes(s, "Notre presentation se deroulera en 7 parties. "
         "Nous commencerons par le probleme concret que nous avons identifie au Cameroun, "
         "puis nous presenterons notre solution, l'architecture, les choix technologiques, "
         "les fonctionnalites cles, l'infrastructure de deploiement, "
         "et nous terminerons par une demonstration en direct et un bilan du projet. "
         "Comptez environ 25 a 30 minutes.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "01 — Problematique")
txb(s, "Verifier un diplome aujourd'hui :\nun processus long, manuel et sans tracabilite",
    Inches(0.55), Inches(0.65), Inches(12), Inches(1.3),
    size=26, color=WHITE, bold=True)

# 6 contextes ou la verification est necessaire
contextes = [
    ("briefcase",   "Recrutement",            "Un employeur veut\nconfirmer le diplome\nd'un candidat"),
    ("award",       "Admission universitaire", "Une universite verifie\nle niveau d'un etudiant\nqui postule en master"),
    ("trending-up", "Bourse d'etudes",         "Un organisme examine\nle parcours academique\navant attribution"),
    ("globe",       "Demande de visa",          "Une ambassade controle\nles diplomes d'un\ncandidant a l'immigration"),
    ("shield",      "Ordre professionnel",      "Medecin, ingenieur, avocat :\nl'ordre verifie le diplome\navant toute admission"),
    ("users",       "Concours & Fonct. pub.",  "Les jurys et commissions\nverifient les titres\ndes candidats"),
]
for i, (ico, title, desc) in enumerate(contextes):
    col = i % 3
    row = i // 3
    x = Inches(0.42 + col * 4.3)
    y = Inches(2.1 + row * 1.65)
    rect(s, NAVY_MED, x, y, Inches(4.05), Inches(1.45))
    rect(s, GOLD, x, y, Inches(4.05), Inches(0.05))
    icon(s, ico, x + Inches(0.15), y + Inches(0.12), size=Inches(0.52))
    txb(s, title, x + Inches(0.82), y + Inches(0.12), Inches(3.1), Inches(0.42),
        size=13, color=GOLD, bold=True)
    txb(s, desc, x + Inches(0.82), y + Inches(0.54), Inches(3.1), Inches(0.82),
        size=11, color=LGRAY, wrap=True)

# Probleme commun a tous ces contextes
rect(s, GOLD, Inches(0.42), Inches(5.5), Inches(12.5), Inches(0.05))
rect(s, NAVY_MED, Inches(0.42), Inches(5.58), Inches(12.5), Inches(1.55))

txb(s, "Dans tous ces cas, le processus est identique :",
    Inches(0.62), Inches(5.65), Inches(12), Inches(0.35),
    size=13, color=GOLD, bold=True)

problemes = [
    ("smartphone", "Appel telephonique au secretariat — qui peut ne pas repondre"),
    ("clock",      "Courrier physique — plusieurs semaines d'attente"),
    ("x-circle",   "Aucun systeme de verification en ligne disponible"),
    ("alert",      "Aucune tracabilite : on ne sait pas qui a verifie quoi, ni quand"),
]
for i, (ico, txt) in enumerate(problemes):
    icon(s, ico, Inches(0.55), Inches(6.1 + i * 0.36), size=Inches(0.28))
    txb(s, txt, Inches(0.95), Inches(6.12 + i * 0.36), Inches(11.8), Inches(0.3),
        size=12, color=LGRAY)

num_badge(s, 3)
notes(s, "La problematique est plus large que la simple fraude a l'embauche. "
         "La verification de diplome concerne de nombreux contextes : "
         "le recrutement, l'admission en master ou en doctorat, les bourses d'etudes, "
         "les demandes de visa, les ordres professionnels comme les medecins ou les ingenieurs, "
         "et les concours de la fonction publique. "
         "Dans tous ces cas, le probleme est le meme : le processus est manuel, lent, et sans tracabilite. "
         "On appelle le secretariat de l'universite — qui peut ne pas repondre. "
         "On envoie un courrier — qui prend plusieurs semaines. "
         "Et a la fin, on n'a aucune preuve horodatee de la verification effectuee.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "02 — Solution")
txb(s, "INUBIL Verify — La solution",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

rect(s, LGRAY, Inches(0.38), Inches(1.5), Inches(12.55), Inches(0.95))
rect(s, GOLD, Inches(0.38), Inches(1.5), Inches(0.06), Inches(0.95))
txb(s, "Une plateforme web qui permet a ISTAMA INUBIL d'emettre des diplomes "
       "certifies sur la blockchain Polygon, et a n'importe qui dans le monde "
       "de verifier leur authenticite en quelques secondes — sans appel, sans courrier.",
    Inches(0.58), Inches(1.62), Inches(12.1), Inches(0.7),
    size=14, color=NAVY)

# 3 piliers avec icones reelles
piliers = [
    ("lock",         "Emission\nsecurisee",
     "Workflow multi-roles :\nagent de saisie → directeur\n→ ancrage blockchain",
     RGBColor(0x22, 0x5C, 0x8A)),
    ("link",         "Certification\nblockchain",
     "Chaque diplome est enregistre\nsur Polygon — immuable\net verifiable publiquement",
     RGBColor(0x83, 0x47, 0xE5)),
    ("check-circle", "Verification\ninstantanee",
     "3 methodes : scan QR code,\nhash SHA-256, ou upload\ndu fichier PDF",
     GREEN),
]
for i, (ico, title, desc, accent) in enumerate(piliers):
    x = Inches(0.38 + i * 4.35)
    rect(s, NAVY, x, Inches(2.65), Inches(4.1), Inches(4.05))
    rect(s, accent, x, Inches(2.65), Inches(4.1), Inches(0.06))
    icon(s, ico, x + Inches(0.25), Inches(2.85), size=Inches(0.95))
    txb(s, title, x + Inches(1.4), Inches(2.8), Inches(2.5), Inches(0.85),
        size=18, color=GOLD, bold=True)
    txb(s, desc, x + Inches(0.2), Inches(3.85), Inches(3.7), Inches(2.65),
        size=13, color=LGRAY)

num_badge(s, 4)
notes(s, "Notre reponse a ce probleme est INUBIL Verify. "
         "L'idee centrale est simple : quand un diplome est valide par le directeur, "
         "son empreinte numerique est immediatement ancree sur la blockchain Polygon. "
         "A partir de ce moment, n'importe qui — recruteur, institution partenaire, ambassade — "
         "peut verifier son authenticite en quelques secondes, depuis n'importe ou, "
         "sans avoir besoin de contacter l'universite. "
         "La plateforme repose sur 3 piliers : une emission securisee avec separation des roles, "
         "une certification blockchain immuable, et une verification publique instantanee.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "03 — Architecture")
txb(s, "Architecture du systeme",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

# Frontend
rect(s, WHITE, Inches(0.4), Inches(1.48), Inches(3.45), Inches(1.0))
rect(s, GOLD, Inches(0.4), Inches(1.48), Inches(3.45), Inches(0.06))
icon(s, "server", Inches(0.5), Inches(1.6), size=Inches(0.55))
txb(s, "FRONTEND", Inches(1.15), Inches(1.57), Inches(2.6), Inches(0.38),
    size=13, color=NAVY, bold=True)
txb(s, "Angular 17\nDashboard Admin - Espace Etudiant\nVerification Publique",
    Inches(0.5), Inches(2.04), Inches(3.2), Inches(0.55), size=10, color=DGRAY)

txb(s, "->", Inches(3.93), Inches(1.75), Inches(0.45), Inches(0.45),
    size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# Backend
rect(s, NAVY, Inches(4.5), Inches(1.48), Inches(4.3), Inches(1.0))
rect(s, GOLD, Inches(4.5), Inches(1.48), Inches(4.3), Inches(0.06))
icon(s, "cpu", Inches(4.6), Inches(1.58), size=Inches(0.55))
txb(s, "BACKEND - NestJS API REST",
    Inches(5.25), Inches(1.57), Inches(3.4), Inches(0.38),
    size=12, color=GOLD, bold=True)
txb(s, "Auth JWT - RBAC Granulaire\nGestion Diplomes - Multi-tenant\nSwagger OpenAPI",
    Inches(4.6), Inches(1.98), Inches(4.05), Inches(0.55), size=10, color=LGRAY)

for xi in [Inches(5.05), Inches(6.55), Inches(8.05)]:
    txb(s, "v", xi, Inches(2.62), Inches(0.4), Inches(0.35),
        size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# 3 services avec vrais logos
services = [
    ("database", "PostgreSQL", "Base de donnees\nPrisma ORM\nSoft delete", Inches(4.5)),
    ("cloud",    "Cloudflare R2", "Stockage PDF\nPresigned URL\n10 GB gratuit", Inches(6.15)),
    ("link",     "Polygon", "Blockchain\nAmoy Testnet\nSmart Contract", Inches(7.8)),
]
for ico, name, desc, x in services:
    rect(s, WHITE, x, Inches(3.1), Inches(1.6), Inches(1.55))
    rect(s, GOLD, x, Inches(3.1), Inches(1.6), Inches(0.04))
    icon(s, ico, x + Inches(0.08), Inches(3.17), size=Inches(0.45))
    txb(s, name, x + Inches(0.6), Inches(3.18), Inches(0.95), Inches(0.38),
        size=10, color=NAVY, bold=True)
    txb(s, desc, x + Inches(0.08), Inches(3.6), Inches(1.45), Inches(0.95),
        size=9, color=DGRAY)

# Docker
rect(s, NAVY_MED, Inches(0.4), Inches(4.85), Inches(12.55), Inches(0.85))
rect(s, GOLD, Inches(0.4), Inches(4.85), Inches(12.55), Inches(0.04))
tech(s, "docker", Inches(0.5), Inches(4.93), size=Inches(0.55))
txb(s, "Docker Compose  -  Orchestration locale (backend - frontend - PostgreSQL - pgAdmin)",
    Inches(1.2), Inches(4.96), Inches(11.5), Inches(0.38),
    size=12, color=WHITE, bold=True)
txb(s, "Variables d'env via .env  -  Volumes persistants  -  Hot-reload en developpement",
    Inches(1.2), Inches(5.32), Inches(11.5), Inches(0.38),
    size=11, color=LGRAY)

num_badge(s, 5)
notes(s, "L'architecture repose sur 3 couches. "
         "Un frontend Angular pour les interfaces utilisateurs, qui n'est pas encore developpe. "
         "Un backend NestJS qui est le coeur du systeme — il gere toute la logique metier. "
         "Et 3 services externes : PostgreSQL pour les donnees, Cloudflare R2 pour les PDFs, "
         "et Polygon pour la blockchain. "
         "Tout est orchestre avec Docker Compose, ce qui garantit que l'environnement est "
         "identique en developpement et en production.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STACK
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "04 — Stack")
txb(s, "Technologies choisies — et pourquoi",
    Inches(0.5), Inches(0.35), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

LOGO = TECHLOGO

def logo_card(slide, logo_file, name, reason, accent, x, y, w=Inches(6.1), h=Inches(1.35)):
    rect(slide, WHITE, x, y, w, h)
    rect(slide, accent, x, y, Inches(0.06), h)
    lp = os.path.join(LOGO, logo_file)
    if os.path.exists(lp):
        try: slide.shapes.add_picture(lp, x + Inches(0.18), y + Inches(0.22), height=Inches(0.85))
        except: pass
    txb(slide, name, x + Inches(1.25), y + Inches(0.1), w - Inches(1.35), Inches(0.42),
        size=14, color=NAVY, bold=True)
    txb(slide, reason, x + Inches(1.25), y + Inches(0.52), w - Inches(1.35), Inches(0.72),
        size=11, color=DGRAY, wrap=True)

logo_card(s, 'nestjs.png',    "NestJS + TypeScript",
    "Architecture modulaire, typage strict.\nLe compilateur detecte les bugs avant l'execution.",
    RGBColor(0xE0, 0x23, 0x4E), Inches(0.35), Inches(1.45))
logo_card(s, 'postgresql.png',"PostgreSQL 16",
    "Transactions ACID, UUID natifs, soft-delete.\nLa coherence des diplomes est garantie par la base.",
    RGBColor(0x33, 0x6A, 0x99), Inches(0.35), Inches(2.9))
logo_card(s, 'polygon.png',   "Polygon (Blockchain)",
    "Reseau EVM, frais quasi-nuls.\nChaque diplome valide est ancre de facon immuable.",
    RGBColor(0x83, 0x47, 0xE5), Inches(0.35), Inches(4.35))
logo_card(s, 'docker.png',    "Docker Compose",
    "Meme environnement dev et production.\nZero surprise lors du deploiement.",
    RGBColor(0x00, 0x91, 0xE2), Inches(0.35), Inches(5.8))

logo_card(s, 'angular.png',   "Angular 17",
    "Framework structure pour SPA complexes.\nGuards de routes et modules lazy-loaded.",
    RGBColor(0xDD, 0x00, 0x31), Inches(6.85), Inches(1.45))
logo_card(s, 'prisma.png',    "Prisma ORM",
    "Schema type, migrations versionnees.\nAucune requete SQL manuelle = aucune injection.",
    RGBColor(0x2D, 0x3A, 0x4A), Inches(6.85), Inches(2.9))
logo_card(s, 'cloudflare.png',"Cloudflare R2",
    "10 GB gratuits, 0 frais d'egress.\nAPI S3-compatible : migration vers AWS = 0 ligne de code.",
    RGBColor(0xF4, 0x81, 0x20), Inches(6.85), Inches(4.35))
logo_card(s, 'typescript.png',"JWT (Auth stateless)",
    "Access token 15 min + refresh 7 jours.\nSessions hachees SHA-256, jamais en clair.",
    RGBColor(0x31, 0x78, 0xC6), Inches(6.85), Inches(5.8))

num_badge(s, 6)
notes(s, "Chaque choix technologique a une raison precise. "
         "NestJS pour sa structure modulaire qui force la separation des responsabilites. "
         "PostgreSQL pour sa robustesse et ses garanties ACID. "
         "Polygon pour ses couts quasi-nuls par rapport a Ethereum mainnet. "
         "Cloudflare R2 pour le stockage des PDFs — gratuit jusqu'a 10 gigaoctets "
         "et compatible avec l'API AWS S3, ce qui rend la migration triviale en production. "
         "Et Docker pour que l'environnement soit identique du developpement a la production.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WORKFLOW D'EMISSION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "05 — Fonctionnalites")
txb(s, "Workflow d'emission d'un diplome",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

steps = [
    ("01", "Saisie\n(Agent)",         "Creation du brouillon :\ninfos etudiant,\ntype de diplome"),
    ("02", "Upload PDF\n(Agent)",      "Depot du PDF\nCalcul hash SHA-256\nStockage R2"),
    ("03", "Validation\n(Directeur)",  "Verification\nApprobation\nStatut -> actif"),
    ("04", "Blockchain\n(Auto)",       "Enregistrement\nPolygon\nTx confirmee"),
    ("05", "Verifiable\n(Public)",     "Lien unique\nQR code\nVerification"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.38 + i * 2.52)
    step_box(s, num, title, desc, x, Inches(1.5), active=(i == 2))
    if i < 4:
        txb(s, "->", Inches(2.69 + i * 2.52), Inches(2.3), Inches(0.38), Inches(0.45),
            size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# Icones roles
roles = [
    ("users",     "Role : Agent de saisie",   Inches(0.38)),
    ("users",     "Role : Agent de saisie",   Inches(2.9)),
    ("lock",      "Role : Directeur",         Inches(5.42)),
    ("cpu",       "Systeme automatique",      Inches(7.94)),
    ("globe",     "Acces public",             Inches(10.46)),
]
for ico, label, x in roles:
    icon(s, ico, x + Inches(0.75), Inches(3.72), size=Inches(0.38))
    txb(s, label, x, Inches(4.18), Inches(2.35), Inches(0.35),
        size=9, color=DGRAY, align=PP_ALIGN.CENTER)

h_rule(s, Inches(0.38), Inches(4.65), Inches(12.55))

rect(s, WHITE, Inches(0.38), Inches(4.8), Inches(12.55), Inches(1.9))
rect(s, GOLD, Inches(0.38), Inches(4.8), Inches(0.06), Inches(1.9))
txb(s, "Points cles de securite",
    Inches(0.58), Inches(4.88), Inches(5), Inches(0.38), size=13, color=NAVY, bold=True)
garanties = [
    ("check", "Detection de doublons par hash SHA-256 : impossible d'importer deux fois le meme PDF"),
    ("check", "Validation (etape 3) bloquee si aucun PDF n'a ete uploade au prealable"),
    ("check", "Revocation propagee on-chain : meme si la base de donnees est compromise"),
]
for i, (ico, txt) in enumerate(garanties):
    icon(s, ico, Inches(0.58), Inches(5.35 + i * 0.45), size=Inches(0.3))
    txb(s, txt, Inches(0.98), Inches(5.38 + i * 0.45), Inches(11.85), Inches(0.38),
        size=12, color=DGRAY)

num_badge(s, 7)
notes(s, "Voici comment un diplome est emis dans INUBIL Verify. "
         "L'agent de saisie cree le dossier et uploade le PDF du diplome. "
         "Le systeme calcule automatiquement le hash SHA-256 du fichier — une empreinte numerique unique. "
         "Le directeur examine et valide le dossier. C'est ce moment qui declenche l'ancrage blockchain. "
         "La separation entre agent et directeur est un controle qualite voulu : "
         "aucune personne seule ne peut emettre un diplome sans verification. "
         "Si le meme PDF est uploade deux fois, le systeme le detecte et refuse.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SECURITE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "05 — Securite")
txb(s, "Securite — Architecture en couches",
    Inches(0.5), Inches(0.7), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

layers = [
    ("globe",  "Couche 1 — Reseau",      "Helmet (headers HTTP)  -  CORS strict  -  Rate limiting (Throttler)  -  HTTPS"),
    ("key",    "Couche 2 — Auth",         "JWT 15 min  -  Refresh token 7j  -  Sessions hachees SHA-256  -  Blocage 5 tentatives / 15 min"),
    ("shield", "Couche 3 — RBAC",         "25 permissions granulaires  -  Guards NestJS  -  Verification par endpoint"),
    ("layers", "Couche 4 — Multi-tenant", "Isolation totale par universite  -  ForbiddenException cross-tenant  -  Super admin = seul a tout voir"),
    ("eye",    "Couche 5 — Donnees",      "Soft delete  -  Audit trail complet  -  bcrypt salt=12  -  Jamais de secrets en clair"),
]
for i, (ico, layer, desc) in enumerate(layers):
    y = Inches(1.65 + i * 1.0)
    rect(s, NAVY_MED, Inches(0.38), y, Inches(12.55), Inches(0.85))
    rect(s, GOLD, Inches(0.38), y, Inches(0.06), Inches(0.85))
    rect(s, GOLD, Inches(0.38), y, Inches(2.75), Inches(0.85))
    icon(s, ico, Inches(0.48), y + Inches(0.17), size=Inches(0.48))
    txb(s, layer, Inches(1.05), y + Inches(0.2), Inches(1.6), Inches(0.45),
        size=11, color=NAVY, bold=True)
    txb(s, desc, Inches(3.2), y + Inches(0.22), Inches(9.65), Inches(0.45),
        size=12, color=LGRAY)

num_badge(s, 8)
notes(s, "La securite est construite en 5 couches independantes, selon le principe de defense en profondeur. "
         "Si une couche est contournee, les suivantes protegent toujours. "
         "La couche la plus importante pour ce projet est la couche 4 : l'isolation multi-tenant. "
         "Chaque universite ne peut voir que ses propres donnees. "
         "Cette isolation est verifiee sur chaque requete, pas seulement au moment de la connexion. "
         "La couche 5 garantit que toutes les actions sont auditees : on sait qui a fait quoi et quand.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BLOCKCHAIN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "05 — Blockchain")
txb(s, "Integration blockchain Polygon",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

rect(s, LGRAY, Inches(0.38), Inches(1.48), Inches(7.8), Inches(4.65))
rect(s, NAVY, Inches(0.38), Inches(1.48), Inches(7.8), Inches(0.52))
txb(s, "Comment ca fonctionne",
    Inches(0.55), Inches(1.56), Inches(7.4), Inches(0.35),
    size=14, color=WHITE, bold=True)

flow = [
    ("award",       "Diplome valide en base",              True),
    ("code",        "BlockchainService appelle le contrat", False),
    ("lock",        "Transaction signee avec cle privee",  False),
    ("link",        "Confirmee sur Polygon Amoy",          False),
    ("database",    "tx_hash sauvegarde en base",          False),
    ("check-circle","Verifiable sur PolygonScan",          True),
]
for i, (ico, text, highlight) in enumerate(flow):
    y = Inches(2.12 + i * 0.62)
    icon(s, ico, Inches(0.5), y, size=Inches(0.38))
    rect(s, GOLD if highlight else NAVY, Inches(0.95), y + Inches(0.08),
         Inches(0.06), Inches(0.38))
    txb(s, text, Inches(1.12), y + Inches(0.05), Inches(6.8), Inches(0.38),
        size=13, color=NAVY if highlight else DGRAY, bold=highlight)
    if i < 5:
        txb(s, "|", Inches(0.95), y + Inches(0.45), Inches(0.1), Inches(0.2),
            size=9, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, NAVY, Inches(8.4), Inches(1.48), Inches(4.55), Inches(4.65))
rect(s, GOLD, Inches(8.4), Inches(1.48), Inches(4.55), Inches(0.06))
txb(s, "Contrat deploye", Inches(8.58), Inches(1.6), Inches(4.2), Inches(0.38),
    size=13, color=GOLD, bold=True)

infos = [
    ("Reseau",   "Polygon Amoy Testnet"),
    ("Contrat",  "0x75a23a...1315E"),
    ("Langage",  "Solidity 0.8.20"),
    ("Lib",      "OpenZeppelin Ownable"),
    ("",""),
    ("Fonction", "enregistrerDiplome()"),
    ("",         "revoquerDiplome()"),
    ("",         "verifierDiplome()"),
    ("",""),
    ("Pattern",  "Fire & forget"),
    ("",         "=> ne bloque pas l'API"),
    ("",""),
    ("Preuve",   "PolygonScan public"),
]
y_i = Inches(2.05)
for label, val in infos:
    if label:
        txb(s, label + " :", Inches(8.58), y_i, Inches(1.3), Inches(0.32),
            size=11, color=GOLD, bold=True)
    txb(s, val, Inches(9.95), y_i, Inches(2.85), Inches(0.32),
        size=11, color=WHITE)
    y_i += Inches(0.3)

num_badge(s, 9)
notes(s, "L'integration blockchain fonctionne de facon transparente pour l'utilisateur. "
         "Quand le directeur valide un diplome, le backend appelle notre contrat intelligent deploye sur Polygon. "
         "Ce contrat stocke le hash SHA-256 du diplome et l'identifiant de l'universite. "
         "Le pattern utilise est 'fire and forget' : l'API repond immediatement au directeur "
         "sans attendre la confirmation blockchain, qui arrive quelques secondes plus tard. "
         "Cette transaction est publique et immuable. "
         "Si la plateforme disparaissait demain, la preuve resterait sur la blockchain pour toujours.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — VERIFICATION PUBLIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "05 — Verification")
txb(s, "Verification publique — 3 methodes",
    Inches(0.5), Inches(0.35), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

methods = [
    ("smartphone", "QR Code / Lien unique", "GET /verify/{INUB-YYYY-XXXX}",
     "Scan du QR code imprime\nsur le diplome physique.\nResultat instantane.\n30 req/min par IP."),
    ("lock",       "Hash SHA-256",           "POST /verify/hash",
     "Soumettre directement\nle hash du fichier.\nPour les integrateurs\ntechniques. 10 req/min."),
    ("file-text",  "Upload PDF",             "POST /verify/upload",
     "Uploader le fichier PDF.\nLe serveur calcule le hash\net compare. Preuve\ncryptographique. 5 req/min."),
]
for i, (ico, title, endpoint, desc) in enumerate(methods):
    x = Inches(0.38 + i * 4.35)
    rect(s, WHITE, x, Inches(1.52), Inches(4.1), Inches(4.3))
    rect(s, NAVY, x, Inches(1.52), Inches(4.1), Inches(1.0))
    icon(s, ico, x + Inches(0.2), Inches(1.62), size=Inches(0.7))
    txb(s, title, x + Inches(1.1), Inches(1.65), Inches(2.85), Inches(0.85),
        size=14, color=GOLD, bold=True)
    rect(s, LGRAY, x + Inches(0.15), Inches(2.57), Inches(3.8), Inches(0.4))
    txb(s, endpoint, x + Inches(0.2), Inches(2.62), Inches(3.7), Inches(0.3),
        size=11, color=NAVY, bold=True)
    txb(s, desc, x + Inches(0.2), Inches(3.08), Inches(3.7), Inches(2.5),
        size=13, color=DGRAY)

rect(s, NAVY, Inches(0.38), Inches(6.0), Inches(12.55), Inches(1.1))
txb(s, "Resultats possibles :", Inches(0.58), Inches(6.08), Inches(2.5), Inches(0.38),
    size=12, color=GOLD, bold=True)
verdicts = [
    ("check-circle", "AUTHENTIQUE", GREEN),
    ("x-circle",     "REVOQUE",     RED),
    ("alert",        "NON TROUVE",  ORANGE),
    ("shield",       "FALSIFIE",    RED),
]
for i, (ico, v, col) in enumerate(verdicts):
    icon(s, ico, Inches(3.2 + i * 2.35), Inches(6.1), size=Inches(0.4))
    txb(s, v, Inches(3.68 + i * 2.35), Inches(6.15), Inches(1.85), Inches(0.38),
        size=13, color=col, bold=True)

num_badge(s, 10)
notes(s, "La verification est la fonctionnalite publique phare de la plateforme. "
         "Elle est accessible sans compte, sans inscription, sans friction. "
         "Le recruteur qui recoit un diplome a 3 options : "
         "scanner le QR code imprime sur le document, "
         "soumettre le hash SHA-256 s'il dispose d'un outil technique, "
         "ou simplement uploader le fichier PDF. "
         "Dans les 3 cas, le systeme calcule ou compare le hash et retourne un resultat immediatement. "
         "Si le diplome a ete revoqu e, le PDF n'est plus accessible — c'est une protection ajoutee.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — API REST
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "05 — API")
txb(s, "Une API REST complete et documentee",
    Inches(0.5), Inches(0.7), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

kpi(s, "modules\nfonctionnels",  "12", "",       Inches(0.38), Inches(1.62))
kpi(s, "endpoints\ndocumentes",  "50+","routes",  Inches(3.53), Inches(1.62))
kpi(s, "droits\ngranulaires",    "25", "permiss.",Inches(6.68), Inches(1.62))
kpi(s, "tests\nautomatiques",    "60+","jest",    Inches(9.83), Inches(1.62))

txb(s, "Chaque action correspond a un endpoint precise, protege par un droit d'acces specifique.",
    Inches(0.5), Inches(3.35), Inches(12.4), Inches(0.45),
    size=14, color=LGRAY)

zones = [
    ("globe",   "Zone Publique\n(sans connexion)",
     "Verifier un diplome\nTelecharger rapport PDF\nScanner QR code",       GREEN),
    ("users",   "Zone Staff\n(connexion requise)",
     "Creer / valider diplomes\nGerer les etudiants\nVoir l'historique",    GOLD),
    ("shield",  "Zone Admin\n(super admin)",
     "Gerer les universites\nGerer les comptes staff\nJournal d'audit",     PURPLE),
]
for i, (ico, zone, items, color) in enumerate(zones):
    x = Inches(0.38 + i * 4.35)
    rect(s, NAVY_MED, x, Inches(3.95), Inches(4.1), Inches(2.75))
    rect(s, color, x, Inches(3.95), Inches(4.1), Inches(0.06))
    icon(s, ico, x + Inches(0.18), Inches(4.1), size=Inches(0.65))
    txb(s, zone, x + Inches(1.0), Inches(4.1), Inches(2.95), Inches(0.7),
        size=14, color=color, bold=True)
    txb(s, items, x + Inches(0.2), Inches(4.85), Inches(3.7), Inches(1.7),
        size=13, color=LGRAY)

txb(s, "Documentation interactive : http://localhost:3000/api/docs",
    Inches(0.38), Inches(6.88), Inches(12.55), Inches(0.3),
    size=11, color=GOLD, align=PP_ALIGN.CENTER)

num_badge(s, 11)
notes(s, "Le backend est decompose en 12 modules qui correspondent aux 3 zones de l'application. "
         "La zone publique est accessible sans connexion — c'est la que se fait la verification de diplomes. "
         "La zone staff necessite une authentification avec les droits appropries. "
         "La zone admin est reservee aux super administrateurs qui gerent les universites elles-memes. "
         "Toute l'API est documentee de facon interactive sur Swagger OpenAPI. "
         "Le developpeur frontend peut tester chaque endpoint directement depuis le navigateur.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — INFRASTRUCTURE & STOCKAGE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "06 — Infrastructure")
txb(s, "Infrastructure & Choix du stockage",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

# Docker + Polygon
for logo_f, ico_f, title, desc, accent, xi in [
    ('docker.png',  'docker',  "Docker Compose",
     "4 services : backend, frontend, PostgreSQL, pgAdmin.\nHot-reload en dev. Volumes persistants.",
     RGBColor(0x00, 0x91, 0xE2), Inches(0.38)),
    ('polygon.png', 'link',    "Polygon Amoy",
     "Testnet EVM. Contrat : 0x75a23a...1315E.\nVerifiable sur amoy.polygonscan.com",
     RGBColor(0x83, 0x47, 0xE5), Inches(6.75)),
]:
    rect(s, LGRAY, xi, Inches(1.52), Inches(6.15), Inches(1.52))
    rect(s, accent, xi, Inches(1.52), Inches(0.06), Inches(1.52))
    lpath = os.path.join(TECHLOGO, logo_f)
    if os.path.exists(lpath):
        try: s.shapes.add_picture(lpath, xi + Inches(0.15), Inches(1.62), height=Inches(0.72))
        except: pass
    txb(s, title, xi + Inches(1.12), Inches(1.6), Inches(4.85), Inches(0.4),
        size=14, color=NAVY, bold=True)
    txb(s, desc, xi + Inches(1.12), Inches(2.0), Inches(4.85), Inches(0.9),
        size=12, color=DGRAY)

# Comparatif stockage
txb(s, "Choix du stockage PDF : Cloudflare R2 vs AWS S3",
    Inches(0.38), Inches(3.22), Inches(12.55), Inches(0.45),
    size=15, color=NAVY, bold=True)
h_rule(s, Inches(0.38), Inches(3.68), Inches(12.55))

# Cloudflare R2
rect(s, LGRAY, Inches(0.38), Inches(3.82), Inches(5.85), Inches(2.82))
rect(s, RGBColor(0xF4, 0x81, 0x20), Inches(0.38), Inches(3.82), Inches(5.85), Inches(0.06))
lp_cf = os.path.join(TECHLOGO, 'cloudflare.png')
if os.path.exists(lp_cf):
    try: s.shapes.add_picture(lp_cf, Inches(0.5), Inches(3.92), height=Inches(0.52))
    except: pass
txb(s, "Cloudflare R2   —   ACTUEL",
    Inches(1.62), Inches(3.92), Inches(4.45), Inches(0.45),
    size=14, color=NAVY, bold=True)
r2_pros = [
    ("check", GREEN, "10 GB/mois gratuits, 0 frais de sortie"),
    ("check", GREEN, "API 100% compatible AWS S3"),
    ("check", GREEN, "Presigned URLs (TTL 15 min) — PDF inaccessible en direct"),
    ("check", GREEN, "Parfait pour le developpement et la phase MVP"),
]
for j, (ico, col, p) in enumerate(r2_pros):
    icon(s, ico, Inches(0.48), Inches(4.48 + j * 0.48), size=Inches(0.32))
    txb(s, p, Inches(0.88), Inches(4.5 + j * 0.48), Inches(5.2), Inches(0.38),
        size=12, color=col)

# AWS S3
rect(s, LGRAY, Inches(6.55), Inches(3.82), Inches(6.45), Inches(2.82))
rect(s, RGBColor(0xFF, 0x99, 0x00), Inches(6.55), Inches(3.82), Inches(6.45), Inches(0.06))
lp_aws = os.path.join(TECHLOGO, 'aws.png')
if os.path.exists(lp_aws):
    try: s.shapes.add_picture(lp_aws, Inches(6.68), Inches(3.92), height=Inches(0.52))
    except: pass
txb(s, "AWS S3   —   OPTION PRODUCTION",
    Inches(8.05), Inches(3.92), Inches(4.75), Inches(0.45),
    size=14, color=NAVY, bold=True)
s3_items = [
    ("check",   GREEN,  "Ecosysteme mature, SLA 99.999%"),
    ("check",   GREEN,  "Integration CloudFront CDN native"),
    ("x-circle",RED,    "Frais de sortie de donnees (egress)"),
    ("zap",     GOLD,   "Migration : 0 ligne de code (meme API S3)"),
]
for j, (ico, col, p) in enumerate(s3_items):
    icon(s, ico, Inches(6.65), Inches(4.48 + j * 0.48), size=Inches(0.32))
    txb(s, p, Inches(7.05), Inches(4.5 + j * 0.48), Inches(5.75), Inches(0.38),
        size=12, color=col)

rect(s, NAVY, Inches(0.38), Inches(6.75), Inches(12.55), Inches(0.42))
txb(s, "Le choix Cloudflare R2 vs AWS S3 est une decision de configuration, "
       "pas de code — l'API est identique.",
    Inches(0.55), Inches(6.82), Inches(12.3), Inches(0.3),
    size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

num_badge(s, 12)
notes(s, "Pour le stockage des fichiers PDF, nous avons fait un choix delibere : Cloudflare R2. "
         "En developpement et en phase initiale, R2 est parfait — gratuit jusqu'a 10 gigaoctets "
         "et sans aucun frais de sortie de donnees. "
         "Si le projet monte en charge et necessite l'ecosysteme AWS, la migration est triviale. "
         "Les deux services utilisent exactement la meme API S3 : il suffit de changer "
         "3 variables d'environnement dans le fichier .env. "
         "Aucune ligne de code a modifier. C'est pourquoi nous appelons ca une decision de configuration.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — QUALITE & TESTS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "06 — Qualite")
txb(s, "Qualite du code & tests automatiques",
    Inches(0.5), Inches(0.7), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

txb(s, "60+ tests unitaires — chaque service metier est teste independamment.",
    Inches(0.5), Inches(1.6), Inches(12.4), Inches(0.45), size=15, color=LGRAY)

garanties_g = [
    ("x-circle", "Ce qu'on ne PEUT PAS faire",
     ["Creer deux fois le meme diplome (meme PDF)",
      "Valider sans avoir uploade un PDF",
      "Voir les diplomes d'une autre universite",
      "Se connecter apres 5 echecs (15 min de blocage)"],
     RED),
    ("check-circle", "Ce qu'on peut TOUJOURS faire",
     ["Verifier n'importe quel diplome sans compte",
      "Telecharger le rapport PDF de verification",
      "Scanner le QR code depuis un mobile",
      "Consulter l'historique complet des actions"],
     GREEN),
    ("code", "Garanties du code",
     ["TypeScript strict — 0 erreur de type",
      "Validation des donnees a chaque entree API",
      "Journal d'audit sur toute modification",
      "Secrets jamais codes en dur (fichier .env)"],
     GOLD),
    ("server", "Environnement maitrise",
     ["Docker : meme config dev et production",
      "Migrations Prisma versionnees et rejouables",
      "Variables d'env documentees dans .env.example",
      "Swagger : contrat API toujours a jour"],
     BLUE),
]
for i, (ico, title, items, color) in enumerate(garanties_g):
    col = i % 2
    row = i // 2
    x = Inches(0.38 + col * 6.5)
    y = Inches(2.2 + row * 2.45)
    rect(s, NAVY_MED, x, y, Inches(6.2), Inches(2.25))
    rect(s, color, x, y, Inches(6.2), Inches(0.06))
    icon(s, ico, x + Inches(0.15), y + Inches(0.12), size=Inches(0.5))
    txb(s, title, x + Inches(0.82), y + Inches(0.15), Inches(5.25), Inches(0.42),
        size=13, color=color, bold=True)
    for j, item in enumerate(items):
        txb(s, "  " + item, x + Inches(0.15), y + Inches(0.68 + j * 0.36),
            Inches(5.9), Inches(0.33), size=11, color=LGRAY)

num_badge(s, 13)
notes(s, "Comment garantit-on la qualite ? D'abord par les tests automatiques : "
         "60 tests qui verifient les comportements critiques du systeme. "
         "Mais surtout par les contraintes architecturales. "
         "Certaines erreurs sont rendues impossibles par conception : "
         "on ne peut pas creer deux fois le meme diplome, on ne peut pas valider sans PDF. "
         "Ces garanties ne dependent pas du comportement des utilisateurs. "
         "Elles sont enforced par le code.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DEMONSTRATION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "07 — Demo")
txb(s, "Demonstration en direct",
    Inches(0.5), Inches(0.7), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

demo = [
    ("key",         "1", "Connexion",             "admin@inubil.com\nAdmin123!",              Inches(0.38)),
    ("file-text",   "2", "Creer un diplome",       "Remplir le formulaire\npour un etudiant", Inches(3.72)),
    ("check-circle","3", "Valider (Directeur)",    "Upload PDF\npuis Validation",             Inches(7.06)),
    ("globe",       "4", "Verif. publique",        "Scanner le QR code\n/verify/INUB-...",    Inches(10.4)),
]
for ico, num, title, desc, x in demo:
    rect(s, NAVY_MED, x, Inches(1.7), Inches(3.12), Inches(4.05))
    rect(s, GOLD, x, Inches(1.7), Inches(3.12), Inches(0.06))
    icon(s, ico, x + Inches(0.15), Inches(1.85), size=Inches(0.65))
    txb(s, num, x + Inches(0.9), Inches(1.85), Inches(0.5), Inches(0.6),
        size=28, color=GOLD, bold=True)
    txb(s, title, x + Inches(0.15), Inches(2.6), Inches(2.85), Inches(0.5),
        size=14, color=WHITE, bold=True)
    txb(s, desc, x + Inches(0.15), Inches(3.15), Inches(2.85), Inches(0.9),
        size=13, color=LGRAY)
    rect(s, NAVY, x + Inches(0.15), Inches(4.18), Inches(2.82), Inches(1.25))
    txb(s, "[ capture d'ecran ]", x + Inches(0.15), Inches(4.58),
        Inches(2.82), Inches(0.45), size=11, color=DGRAY, align=PP_ALIGN.CENTER)

rect(s, GOLD, Inches(0.38), Inches(5.95), Inches(12.55), Inches(0.52))
txb(s, "API : http://localhost:3000   -   Swagger : http://localhost:3000/api/docs   -   App : http://localhost:4200",
    Inches(0.5), Inches(6.03), Inches(12.3), Inches(0.38),
    size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

num_badge(s, 14)
notes(s, "Passons a la demonstration. "
         "Nous allons parcourir le flux complet en 4 etapes. "
         "D'abord la connexion en tant qu'administrateur avec les identifiants du seed. "
         "Ensuite la creation d'un diplome pour un etudiant fictif. "
         "Puis la validation par le directeur apres upload du PDF — c'est ce moment qui declenche la blockchain. "
         "Et enfin la verification depuis la page publique en scannant le QR code genere. "
         "Ceci montre le systeme en condition reelle, de bout en bout.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PREUVE BLOCKCHAIN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "07 — Blockchain")
txb(s, "Preuve blockchain — Transaction reelle",
    Inches(0.5), Inches(0.35), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

rect(s, LGRAY, Inches(0.38), Inches(1.48), Inches(12.55), Inches(1.45))
rect(s, NAVY, Inches(0.38), Inches(1.48), Inches(0.06), Inches(1.45))
icon(s, "link", Inches(0.5), Inches(1.58), size=Inches(0.5))
txb(s, "Contrat deploye",
    Inches(1.12), Inches(1.58), Inches(3), Inches(0.38),
    size=12, color=NAVY, bold=True)
txb(s, "0x75a23a30753548e08fD4b47e6664a941C6C1315E",
    Inches(1.12), Inches(1.96), Inches(8), Inches(0.38),
    size=13, color=DGRAY, bold=True)
txb(s, "Polygon Amoy Testnet  -  amoy.polygonscan.com",
    Inches(1.12), Inches(2.35), Inches(8), Inches(0.38),
    size=11, color=DGRAY, italic=True)

rect(s, LGRAY, Inches(0.38), Inches(3.1), Inches(12.55), Inches(2.3))
rect(s, GREEN, Inches(0.38), Inches(3.1), Inches(0.06), Inches(2.3))
icon(s, "check-circle", Inches(0.5), Inches(3.2), size=Inches(0.5))
txb(s, "Transaction d'enregistrement diplome",
    Inches(1.12), Inches(3.2), Inches(11), Inches(0.38),
    size=14, color=NAVY, bold=True)
txb(s, "Hash tx : 0xd33eb14ef6e765...", Inches(1.12), Inches(3.62), Inches(11), Inches(0.38),
    size=12, color=DGRAY)
txb(s, "Statut : Succes   -   Reseau : Polygon Amoy   -   Bloc confirme",
    Inches(1.12), Inches(4.0), Inches(11), Inches(0.38), size=12, color=GREEN, bold=True)
txb(s, "Contenu : enregistrerDiplome(bytes32 hashSHA256, bytes32 universiteId, string numeroUnique)",
    Inches(1.12), Inches(4.38), Inches(11), Inches(0.38), size=11, color=DGRAY, italic=True)

rect(s, NAVY, Inches(0.38), Inches(5.58), Inches(12.55), Inches(1.45))
rect(s, GOLD, Inches(0.38), Inches(5.58), Inches(12.55), Inches(0.06))
txb(s, "Ce que ca prouve :", Inches(0.58), Inches(5.7), Inches(3), Inches(0.38),
    size=13, color=GOLD, bold=True)
preuves = [
    ("award",  "Le diplome existait a la date de la transaction"),
    ("lock",   "Son hash SHA-256 est immuable sur la blockchain"),
    ("globe",  "N'importe qui peut verifier sur PolygonScan sans compte"),
]
for i, (ico, p) in enumerate(preuves):
    icon(s, ico, Inches(4.0 + i * 3.1), Inches(5.72), size=Inches(0.38))
    txb(s, p, Inches(4.48 + i * 3.1), Inches(5.75), Inches(2.55), Inches(0.65),
        size=12, color=WHITE, wrap=True)

num_badge(s, 15)
notes(s, "Voici la preuve concrete que la blockchain fonctionne. "
         "Cette transaction est reelle, elle est sur le reseau Polygon Amoy. "
         "Elle contient le hash SHA-256 d'un diplome de test. "
         "Cette preuve est publique, permanente, et ne peut pas etre modifiee. "
         "N'importe qui peut la verifier en cherchant l'adresse de notre contrat sur amoy.polygonscan.com. "
         "C'est la garantie ultime : meme si notre serveur tombe, la preuve reste sur la blockchain.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DEFIS & SOLUTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, LGRAY)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "07 — Defis")
txb(s, "Defis rencontres & solutions apportees",
    Inches(0.5), Inches(0.35), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

defis = [
    ("Module ethers.js introuvable",
     "Les node_modules Docker etaient mis en cache.\nSolution : --renew-anon-volumes force la recreation du volume."),
    ("Transaction hash absent dans la verification",
     "lireBlockchain() ne retourne pas le hash.\nSolution : merger doc.transaction_hash depuis la base de donnees."),
    ("Faucet POL testnet inaccessible",
     "4 faucets essayes sans succes.\nSolution : Triangle Platform — sans compte, sans prerequis."),
    ("Separation upload/validation",
     "Le directeur uploadait et validait en meme temps.\nSolution : 2 endpoints separes agent/directeur."),
    ("Cle privee sans prefixe 0x",
     "Hardhat rejetait la cle silencieusement.\nSolution : ajout du prefixe 0x dans le fichier .env."),
    ("PDF accessible apres revocation",
     "getPdfUrl() ne verifiait pas le statut.\nSolution : ForbiddenException si statut === 'revoque'."),
]
for i, (prob, sol) in enumerate(defis):
    col = i % 2
    row = i // 2
    x = Inches(0.38 + col * 6.55)
    y = Inches(1.5 + row * 1.85)
    rect(s, WHITE, x, y, Inches(6.25), Inches(1.7))
    icon(s, "alert",         x + Inches(0.15), y + Inches(0.12), size=Inches(0.42))
    icon(s, "check-circle",  x + Inches(0.15), y + Inches(0.92), size=Inches(0.38))
    rect(s, RED, x, y, Inches(0.06), Inches(0.85))
    rect(s, GREEN, x, y + Inches(0.85), Inches(0.06), Inches(0.85))
    txb(s, prob, x + Inches(0.7), y + Inches(0.1),
        Inches(5.42), Inches(0.65), size=12, color=NAVY, bold=True)
    txb(s, sol, x + Inches(0.7), y + Inches(0.85),
        Inches(5.42), Inches(0.75), size=11, color=DGRAY)

num_badge(s, 16)
notes(s, "Je vais vous citer les defis techniques les plus significatifs. "
         "Le plus inattendu a ete le probleme de cache Docker : ethers.js n'etait pas charge "
         "car un volume Docker conservait d'anciens node_modules sans la nouvelle dependance. "
         "La solution : forcer la recreation du volume avec --renew-anon-volumes. "
         "Un autre defi etait l'acces au reseau de test Polygon : les faucets classiques "
         "demandent des prerequisites. Nous avons trouve Triangle Platform qui n'en demande pas. "
         "Chaque obstacle nous a appris quelque chose sur l'integration de technologie blockchain en production.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — BILAN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

section_pill(s, "07 — Bilan")
txb(s, "Ce qui a ete livre",
    Inches(0.5), Inches(0.7), Inches(10), Inches(0.75),
    size=28, color=WHITE, bold=True)

kpi(s, "endpoints REST\ndocumentes", "50+", "",     Inches(0.38), Inches(1.62))
kpi(s, "permissions\nRBAC",          "25",  "",     Inches(3.53), Inches(1.62))
kpi(s, "tests unitaires\nJest",       "60+", "",     Inches(6.68), Inches(1.62))
kpi(s, "methodes de\nverification",   "3",   "",     Inches(9.83), Inches(1.62))

rect(s, NAVY_MED, Inches(0.38), Inches(3.42), Inches(6.15), Inches(3.25))
rect(s, GOLD, Inches(0.38), Inches(3.42), Inches(6.15), Inches(0.06))
txb(s, "Fonctionnalites implementees",
    Inches(0.55), Inches(3.5), Inches(5.8), Inches(0.38), size=14, color=GOLD, bold=True)
features = [
    ("check", "Authentification complete (JWT, refresh, sessions, 2FA email)"),
    ("check", "Workflow diplome 5 etapes avec ancrage blockchain"),
    ("check", "Verification publique (QR - hash - upload PDF)"),
    ("check", "Rapport PDF horodate avec logo INUBIL"),
    ("check", "Partage securise par token (espace etudiant)"),
    ("check", "Journal d'audit complet"),
    ("check", "Isolation multi-tenant par universite"),
    ("check", "Stockage securise Cloudflare R2"),
]
for i, (ico, f) in enumerate(features):
    icon(s, ico, Inches(0.48), Inches(4.02 + i * 0.31), size=Inches(0.24))
    txb(s, f, Inches(0.82), Inches(4.04 + i * 0.31), Inches(5.6), Inches(0.28),
        size=11, color=LGRAY)

rect(s, NAVY_MED, Inches(6.75), Inches(3.42), Inches(6.2), Inches(3.25))
rect(s, GOLD, Inches(6.75), Inches(3.42), Inches(6.2), Inches(0.06))
txb(s, "Frontend Angular (a venir)",
    Inches(6.92), Inches(3.5), Inches(5.8), Inches(0.38), size=14, color=GOLD, bold=True)
remaining = [
    ("trending-up", "39 pages identifiees dans la documentation"),
    ("server",      "Dashboard admin + espace universite"),
    ("users",       "Portail etudiant + partage diplomes"),
    ("globe",       "Page verification publique (mobile-first)"),
    ("zap",         "Notifications email (config SMTP en place)"),
    ("cloud",       "Deploiement VPS production"),
    ("link",        "Migration Polygon Mainnet"),
]
for i, (ico, r) in enumerate(remaining):
    icon(s, ico, Inches(6.88), Inches(4.02 + i * 0.31), size=Inches(0.24))
    txb(s, r, Inches(7.22), Inches(4.04 + i * 0.31), Inches(5.6), Inches(0.28),
        size=11, color=LGRAY)

num_badge(s, 17)
notes(s, "En termes de livrable, le backend est complet et fonctionnel. "
         "50 endpoints documentes, 25 droits d'acces granulaires, 60 tests automatiques, "
         "une integration blockchain reelle sur Polygon. "
         "Le frontend Angular est la prochaine etape, en cours de developpement. "
         "39 ecrans ont ete identifies et documentes pour guider NGANGUE TSAFACK BELVIE SCINDIE "
         "dans le developpement des interfaces utilisateurs.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
rect(s, NAVY, 0, 0, Inches(0.22), H)
rect(s, NAVY, 0, 0, W, Inches(1.3))

section_pill(s, "07 — Roadmap")
txb(s, "Roadmap — Ce qui vient ensuite",
    Inches(0.5), Inches(0.35), Inches(11), Inches(0.75),
    size=28, color=WHITE, bold=True)

phases = [
    ("server",      "Phase 1", "Frontend Angular", "2-3 mois",
     ["Dashboard admin complet", "Portail etudiant",
      "Page verification publique", "PWA mobile-first"],
     GOLD),
    ("cloud",       "Phase 2", "Production", "1 mois",
     ["VPS Ubuntu + Nginx + SSL", "Polygon Mainnet",
      "Domaine verify.inubil.com", "Backups automatiques PostgreSQL"],
     GREEN),
    ("trending-up", "Phase 3", "Ameliorations", "Continu",
     ["Notifications email (config prete)", "Export CSV journaux audit",
      "Import CSV etudiants en masse", "Extension a d'autres universites"],
     BLUE),
]
for i, (ico, phase, title, duration, items, color) in enumerate(phases):
    x = Inches(0.38 + i * 4.35)
    rect(s, LGRAY, x, Inches(1.5), Inches(4.1), Inches(5.5))
    rect(s, color, x, Inches(1.5), Inches(4.1), Inches(0.06))
    rect(s, NAVY, x, Inches(1.5), Inches(4.1), Inches(1.15))
    icon(s, ico, x + Inches(0.15), Inches(1.62), size=Inches(0.62))
    txb(s, phase, x + Inches(0.92), Inches(1.6), Inches(1.5), Inches(0.38),
        size=12, color=LGRAY)
    txb(s, title, x + Inches(0.15), Inches(2.02), Inches(3.8), Inches(0.5),
        size=16, color=WHITE, bold=True)
    rect(s, color, x + Inches(2.5), Inches(1.65), Inches(1.45), Inches(0.3))
    txb(s, duration, x + Inches(2.55), Inches(1.67), Inches(1.35), Inches(0.28),
        size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        icon(s, "check", x + Inches(0.15), Inches(2.85 + j * 0.7), size=Inches(0.3))
        txb(s, item, x + Inches(0.55), Inches(2.9 + j * 0.7), Inches(3.4), Inches(0.55),
            size=13, color=DGRAY)

num_badge(s, 18)
notes(s, "La prochaine etape immediate est le developpement du frontend Angular. "
         "Ensuite viendra le deploiement en production sur un VPS avec Nginx, SSL, "
         "et la migration du contrat vers le reseau principal Polygon. "
         "A plus long terme, nous prevoyons d'ajouter des notifications email "
         "— la configuration SMTP est deja en place — "
         "des exports de donnees pour l'administration, "
         "et une extension potentielle a d'autres universites partenaires d'INUBIL.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)

txb(s, "En resume", Inches(0.6), Inches(0.5), Inches(11), Inches(0.55),
    size=13, color=GOLD, bold=True)
txb(s, "INUBIL Verify :", Inches(0.6), Inches(1.1), Inches(11), Inches(0.7),
    size=34, color=WHITE, bold=True)

points = [
    ("lock",         "Une API REST robuste",
     "12 modules, 50+ endpoints, securisee par JWT + RBAC + isolation multi-tenant"),
    ("link",         "Un ancrage blockchain reel",
     "Contrat deploye sur Polygon Amoy, verifiable publiquement sur PolygonScan"),
    ("cloud",        "Un stockage securise",
     "PDFs prives sur Cloudflare R2, accessibles uniquement via presigned URLs 15 min"),
    ("check-circle", "3 modes de verification",
     "QR code  -  Hash SHA-256  -  Upload PDF — sans compte, sans friction"),
    ("server",       "Pret pour la production",
     "Docker, variables d'env, migrations Prisma, tests Jest — il reste le frontend"),
]
for i, (ico, title, desc) in enumerate(points):
    y = Inches(2.05 + i * 0.98)
    icon(s, ico, Inches(0.6), y + Inches(0.08), size=Inches(0.62))
    h_rule(s, Inches(1.4), y + Inches(0.12), Inches(0.04), GOLD, Inches(0.58))
    txb(s, title, Inches(1.58), y + Inches(0.08), Inches(3.5), Inches(0.38),
        size=15, color=GOLD, bold=True)
    txb(s, desc, Inches(1.58), y + Inches(0.44), Inches(11.1), Inches(0.38),
        size=13, color=LGRAY)

rect(s, NAVY_MED, Inches(0.6), Inches(7.08), Inches(12.2), Inches(0.35))
txb(s, "\"La blockchain ne remplace pas la confiance — elle la rend inutile.\"",
    Inches(0.8), Inches(7.12), Inches(12.0), Inches(0.28),
    size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

num_badge(s, 19)
notes(s, "Pour conclure, INUBIL Verify resout un probleme concret : "
         "la verification des diplomes au Cameroun est longue, couteuse et peu fiable. "
         "Notre solution est robuste, securisee, et basee sur une technologie immuable. "
         "Le backend est complet et livre aujourd'hui avec une infrastructure prete pour la production. "
         "Le frontend va permettre a l'equipe d'ISTAMA INUBIL de commencer a utiliser la plateforme "
         "et aux etudiants d'acceder a leurs diplomes certifies.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, NAVY)
rect(s, GOLD, 0, 0, Inches(0.22), H)
h_rule(s, Inches(0.22), Inches(3.22), W - Inches(0.22))

img(s, INUBIL_LOGO,  Inches(0.6),  Inches(0.35), h=Inches(1.4))
rect(s, GOLD, Inches(2.55), Inches(0.42), Inches(0.04), Inches(1.1))
img(s, EVRY_LOGO,   Inches(2.75), Inches(0.3),  h=Inches(2.5))

txb(s, "Merci de votre attention.",
    Inches(0.6), Inches(1.6), Inches(11.5), Inches(1.1),
    size=40, color=WHITE, bold=True)

txb(s, "Questions & Discussion",
    Inches(0.6), Inches(3.55), Inches(11.5), Inches(0.85),
    size=30, color=GOLD, bold=True)

contacts = [
    ("server",  "API",     "http://localhost:3000/api/docs"),
    ("link",    "Contrat", "amoy.polygonscan.com  -  0x75a23a...1315E"),
    ("code",    "Depot",   "github.com/BelCodeur2005/INUBIL-VERIFY"),
]
for i, (ico, label, val) in enumerate(contacts):
    icon(s, ico, Inches(0.6), Inches(4.72 + i * 0.52), size=Inches(0.38))
    txb(s, label + " :", Inches(1.1), Inches(4.75 + i * 0.52),
        Inches(1.2), Inches(0.38), size=13, color=GOLD, bold=True)
    txb(s, val, Inches(2.35), Inches(4.75 + i * 0.52),
        Inches(10), Inches(0.38), size=13, color=LGRAY)

h_rule(s, Inches(0.6), Inches(6.18), Inches(12.2))
txb(s, "TCHENTCHEU JIAGAM FLANC BEL  -  Dev Backend",
    Inches(0.6), Inches(6.28), Inches(5.9), Inches(0.35),
    size=11, color=WHITE, bold=True)
txb(s, "NGANGUE TSAFACK BELVIE SCINDIE  -  Dev Frontend & Maquettes",
    Inches(6.7), Inches(6.28), Inches(6.2), Inches(0.35),
    size=11, color=WHITE, bold=True)

img(s, INUBIL_LOGO, Inches(0.6), Inches(6.72), h=Inches(0.65))
img(s, EVRY_LOGO,   Inches(2.1), Inches(6.68), h=Inches(1.1))
txb(s, "ISTAMA INUBIL  -  Douala, Cameroun  -  2026",
    Inches(4.5), Inches(6.85), Inches(8.8), Inches(0.35),
    size=11, color=DGRAY, align=PP_ALIGN.RIGHT)

notes(s, "Merci de votre attention. "
         "Nous sommes disponibles pour repondre a vos questions sur n'importe quel aspect du projet. "
         "Vous pouvez nous interroger sur l'architecture, les choix technologiques, "
         "la securite, la blockchain, ou les prochaines etapes du developpement. "
         "La documentation Swagger est accessible en direct si vous souhaitez voir les endpoints. "
         "Nous avons egalement le code source disponible si vous voulez l'explorer.")

# ══════════════════════════════════════════════════════════════════════════════
# SAUVEGARDE + TRANSITIONS FADE
# ══════════════════════════════════════════════════════════════════════════════
OUT = r"D:\Bel\projets\projet_plteforme_dauthentification_diplome_istama_inubil\INUBIL-VERIFY\INUBIL-VERIFY-Soutenance.pptx"
prs.save(OUT)

import zipfile, os as _os

TRANSITION = '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>'
TMP = OUT.replace('.pptx', '-tmp.pptx')

with zipfile.ZipFile(OUT, 'r') as zin, zipfile.ZipFile(TMP, 'w', zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml') and 'slideLayout' not in item.filename:
            content = data.decode('utf-8')
            if '<p:transition' not in content:
                content = content.replace('</p:sld>', TRANSITION + '</p:sld>')
            data = content.encode('utf-8')
        zout.writestr(item, data)

_os.replace(TMP, OUT)
print("OK :", OUT)
print("Taille :", round(_os.path.getsize(OUT) / 1024), "KB")
